#!/usr/bin/env python3
"""
prism_cli.py
Prism CLI — 独立命令行入口，用于 OpenClaw Skill 调用。
将 .docx 转换为高质量 PPT，无需 Web 服务/数据库。

用法：
  python3 prism_cli.py --input file.docx --style xiaomi-dark --output-dir ./output
  python3 prism_cli.py --input file.docx --style xiaomi-light --api-key sk-xxx --base-url https://api.openai.com/v1
  python3 prism_cli.py --list-styles
  python3 prism_cli.py --check-deps
"""
from __future__ import annotations

import os
import sys
import json
import asyncio
import argparse
import base64
import shutil
import tempfile
import traceback
from pathlib import Path
from datetime import datetime

# 项目根目录
PROJECT_ROOT = Path(__file__).resolve().parent

# 确保项目根目录在 sys.path 中
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))


# ─── 工具函数 ──────────────────────────────────────────────────────────────────

def log(msg: str, level: str = "info"):
    """带时间戳的日志输出"""
    ts = datetime.now().strftime("%H:%M:%S")
    prefix = {"info": "[*]", "ok": "[+]", "err": "[!]", "warn": "[~]"}.get(level, "[*]")
    print(f"  {prefix} {ts} {msg}")


def list_styles() -> list[dict]:
    """列出所有可用风格"""
    from backend.core.paths import CONFIG_DIR
    styles_dir = CONFIG_DIR / "styles"
    styles = []
    for f in sorted(styles_dir.glob("*.json")):
        data = json.loads(f.read_text(encoding="utf-8"))
        styles.append({
            "id": f.stem,
            "name": data.get("name", f.stem),
            "mixed": data.get("mixed", False),
        })
    return styles


def check_dependencies() -> dict:
    """检查所有必要依赖"""
    results = {}

    # Python 包
    for pkg in ["docx", "pptx", "PIL", "openai", "lxml", "playwright"]:
        try:
            __import__(pkg)
            results[pkg] = True
        except ImportError:
            results[pkg] = False

    # Playwright 浏览器
    try:
        from playwright.sync_api import sync_playwright
        with sync_playwright() as p:
            browser = p.chromium.launch()
            browser.close()
        results["playwright_chromium"] = True
    except Exception:
        results["playwright_chromium"] = False

    # 字体文件
    from backend.core.paths import ASSETS_DIR
    font_dir = ASSETS_DIR / "fonts"
    results["fonts"] = font_dir.exists() and any(font_dir.glob("*.ttf"))

    return results


def save_json(path: str | Path, data: dict | list):
    """将 JSON 数据写入文件。"""
    path = Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")


def load_style(style_name: str) -> dict:
    """加载风格配置。"""
    from backend.core.paths import CONFIG_DIR

    style_path = CONFIG_DIR / "styles" / f"{style_name}.json"
    if not style_path.exists():
        available = [s["id"] for s in list_styles()]
        raise FileNotFoundError(
            f"风格 '{style_name}' 不存在。可用风格: {', '.join(available)}"
        )
    return json.loads(style_path.read_text(encoding="utf-8"))


def profile_pages_data(pages_data: dict) -> dict:
    """为 pages.json 补充内容画像，供 host/custom 两种模式复用。"""
    from backend.core.content_profiler import generate_profile

    pages = pages_data.get("pages", [])
    total_pages = len(pages)
    for i, page in enumerate(pages):
        if not page.get("content_profile"):
            page["content_profile"] = generate_profile(
                page,
                is_first=(i == 0),
                is_last=(i == total_pages - 1),
            )
    pages_data["total_pages"] = total_pages
    return pages_data


def prepare_docx_workdir(docx_path: str, work_dir: str) -> dict:
    """解析 docx 并写入 host-mode 所需的中间文件。"""
    from backend.core.docx_parser import parse_docx

    os.makedirs(work_dir, exist_ok=True)
    pages_data = parse_docx(docx_path, work_dir)
    pages_data = profile_pages_data(pages_data)
    save_json(Path(work_dir) / "pages.json", pages_data)
    return pages_data


def load_pages_data(work_dir: str) -> dict:
    """读取 work_dir 下的 pages.json。"""
    pages_path = Path(work_dir) / "pages.json"
    if not pages_path.exists():
        raise FileNotFoundError(f"缺少 pages.json: {pages_path}")
    pages_data = json.loads(pages_path.read_text(encoding="utf-8"))
    if pages_data.get("pages") and not pages_data["pages"][0].get("content_profile"):
        pages_data = profile_pages_data(pages_data)
        save_json(pages_path, pages_data)
    return pages_data


def load_layout_intents(intents_path: str | None, pages: list, style: dict) -> list:
    """加载布局意图；缺失时使用规则兜底。"""
    from backend.core.layout_planner import _fallback_intent

    if intents_path:
        raw_path = Path(intents_path)
        raw = raw_path.read_text(encoding="utf-8")
        data = json.loads(raw)
        if isinstance(data, dict):
            result = []
            for page in pages:
                entry = {
                    "page": page["page_number"],
                    "layout_intent": data.get(str(page["page_number"])) or data.get(page["page_number"]),
                }
                if not entry["layout_intent"]:
                    entry["layout_intent"] = _fallback_intent(page)
                if style.get("mixed"):
                    entry["contrast_affinity"] = 0
                result.append(entry)
            return result
        if isinstance(data, list):
            return data
        raise ValueError("intents.json 必须是数组或页码映射对象")

    result = []
    for page in pages:
        entry = {"page": page["page_number"], "layout_intent": _fallback_intent(page)}
        if style.get("mixed"):
            entry["contrast_affinity"] = 0
        result.append(entry)
    return result


def emit_layout_plan_prompt_bundle(work_dir: str, style_name: str, prompt_output: str | None = None) -> dict:
    """输出布局规划 prompt，供宿主模型推理。"""
    from backend.core.layout_planner import build_layout_plan_prompts

    pages_data = load_pages_data(work_dir)
    style = load_style(style_name)
    system_prompt, user_prompt = build_layout_plan_prompts(pages_data["pages"], style)
    bundle = {
        "mode": "host_layout_planner",
        "style": style_name,
        "pages": len(pages_data["pages"]),
        "system_prompt": system_prompt,
        "user_prompt": user_prompt,
        "output_format": (
            '[{"page":1,"layout_intent":"..."'
            + (',"contrast_affinity":0' if style.get("mixed") else "")
            + "}]"
        ),
    }
    if prompt_output:
        save_json(prompt_output, bundle)
    return bundle


def emit_page_prompt_bundle(
    work_dir: str,
    style_name: str,
    page_number: int,
    intents_path: str | None = None,
    prompt_output: str | None = None,
) -> dict:
    """输出单页 HTML 生成 prompt，供宿主模型逐页推理。"""
    from backend.core.pipeline import _build_page_style, _rhythm_engine
    from backend.core.prompt_builder import build_system_prompt, build_user_prompt

    pages_data = load_pages_data(work_dir)
    pages = pages_data["pages"]
    total_pages = len(pages)
    if page_number < 1 or page_number > total_pages:
        raise ValueError(f"页码超出范围: {page_number} / {total_pages}")

    style = load_style(style_name)
    intents = load_layout_intents(intents_path, pages, style)
    intent_map = {item["page"]: item["layout_intent"] for item in intents}
    if style.get("mixed"):
        page_tones = _rhythm_engine(total_pages, style, intents)
    else:
        page_tones = ["default"] * total_pages

    prompts_cfg = json.loads((PROJECT_ROOT / "config" / "prompts.json").read_text(encoding="utf-8"))
    page = pages[page_number - 1]
    page_style = _build_page_style(style, page_tones[page_number - 1])
    layout_intent = intent_map.get(page_number, "文字居中排布，标题大字，正文适中，适量留白")

    bundle = {
        "mode": "host_html_generator",
        "style": style_name,
        "page": page_number,
        "page_tone": page_tones[page_number - 1],
        "system_prompt": build_system_prompt(prompts_cfg),
        "user_prompt": build_user_prompt(page, layout_intent, page_style, total_pages),
        "output_path": str(Path(work_dir) / "sections" / f"page_{page_number}.html"),
        "output_format": f'只输出一个 <section class="page-{page_number}"> ... </section>',
    }
    if prompt_output:
        save_json(prompt_output, bundle)
    return bundle


def assemble_html_from_workdir(work_dir: str, style_name: str, html_output: str | None = None) -> str:
    """将 work_dir 下的分段 HTML 组装为完整 HTML。"""
    from backend.core.html_builder import assemble_html

    pages_data = load_pages_data(work_dir)
    total_pages = len(pages_data["pages"])
    sections_dir = Path(work_dir) / "sections"
    missing = []
    sections = []
    for i in range(1, total_pages + 1):
        section_path = sections_dir / f"page_{i}.html"
        if not section_path.exists():
            missing.append(str(section_path))
            continue
        sections.append(section_path.read_text(encoding="utf-8"))
    if missing:
        raise FileNotFoundError("缺少页面 HTML:\n" + "\n".join(missing))

    images_dir = Path(work_dir) / "images"
    images_b64 = {}
    if images_dir.exists():
        for f in images_dir.iterdir():
            if f.suffix.lower() in {".jpg", ".jpeg", ".png"}:
                images_b64[f.stem] = base64.b64encode(f.read_bytes()).decode()

    style = load_style(style_name)
    full_html = assemble_html(sections, images_b64, style)
    output_path = Path(html_output) if html_output else Path(work_dir) / "output.html"
    output_path.write_text(full_html, encoding="utf-8")
    return str(output_path)


def export_from_workdir(work_dir: str, output_dir: str, style_name: str, export_modes: list[str]) -> list[str]:
    """从 work_dir 中间产物导出最终 PPT。"""
    os.makedirs(output_dir, exist_ok=True)
    html_path = Path(work_dir) / "output.html"
    if not html_path.exists():
        raise FileNotFoundError(f"缺少 output.html: {html_path}")

    style = load_style(style_name)
    output_files = []

    if "image" in export_modes:
        pptx_image_path = os.path.join(output_dir, "presentation_image.pptx")
        _build_image_pptx_standalone(str(html_path), pptx_image_path)
        output_files.append(pptx_image_path)

    if "editable" in export_modes:
        from backend.core.layout_extractor import extract_layout
        from backend.core.pptx_builder import build_pptx

        layout_data = extract_layout(str(html_path))
        pptx_edit_path = os.path.join(output_dir, "presentation_editable.pptx")
        build_pptx(layout_data, style, pptx_edit_path, str(Path(work_dir) / "images"))
        output_files.append(pptx_edit_path)

    return output_files


# ─── 核心 Pipeline（无数据库依赖）──────────────────────────────────────────────

async def run_cli_pipeline(
    docx_path: str,
    style_name: str,
    output_dir: str,
    export_modes: list[str] = None,
):
    """
    完整的 CLI pipeline：解析 DOCX → AI 生成 HTML → 导出 PPTX。
    不依赖 database.py / routes.py。
    """
    if export_modes is None:
        export_modes = ["image", "editable"]

    os.makedirs(output_dir, exist_ok=True)

    # ── 准备工作目录 ──
    work_dir = tempfile.mkdtemp(prefix="prism_")
    log(f"工作目录: {work_dir}")

    try:
        # ── Phase 1: 解析 DOCX ──
        log("Phase 1: 解析 DOCX 文件...")
        pages_data = prepare_docx_workdir(docx_path, work_dir)
        pages = pages_data.get("pages", [])
        total_pages = len(pages)

        if total_pages == 0:
            log("DOCX 中未找到任何页面（需要 H1 标题分页）", "err")
            return False

        log(f"解析完成：共 {total_pages} 页", "ok")

        # ── 加载风格 ──
        try:
            style = load_style(style_name)
        except FileNotFoundError as e:
            log(str(e), "err")
            return False
        log(f"使用风格: {style.get('name', style_name)}")

        # ── Phase 2: 内容分析 + 布局规划 ──
        log("Phase 2: 分析内容...")
        log("Phase 2: AI 规划布局...")
        from backend.core.layout_planner import plan_layout
        intents = await asyncio.to_thread(plan_layout, pages, style)
        intent_map = {item["page"]: item["layout_intent"] for item in intents}
        log("布局规划完成", "ok")

        # 色调分配（混合模式）
        from backend.core.pipeline import _rhythm_engine, _build_page_style
        if style.get("mixed"):
            page_tones = _rhythm_engine(total_pages, style, intents)
        else:
            page_tones = ["default"] * total_pages

        # ── Phase 3: 逐页生成 HTML ──
        log(f"Phase 3: 并发生成 {total_pages} 页 HTML...")
        from backend.core.prompt_builder import build_system_prompt, build_user_prompt
        from backend.core.llm_client import call_llm

        prompts_cfg = json.loads((CONFIG_DIR / "prompts.json").read_text(encoding="utf-8"))
        semaphore = asyncio.Semaphore(3)

        async def gen_page(idx: int) -> dict:
            async with semaphore:
                page = pages[idx]
                intent = intent_map.get(page["page_number"],
                                        "文字居中排布，标题大字，正文适中，适量留白")
                page_style = _build_page_style(style, page_tones[idx])
                sys_prompt = build_system_prompt(prompts_cfg)
                user_prompt = build_user_prompt(page, intent, page_style, total_pages)

                for attempt in range(3):
                    try:
                        log(f"  第 {idx+1}/{total_pages} 页生成中...")
                        html = await call_llm(sys_prompt, user_prompt)
                        # 保存单页 HTML
                        sections_dir = os.path.join(work_dir, "sections")
                        os.makedirs(sections_dir, exist_ok=True)
                        with open(os.path.join(sections_dir, f"page_{idx+1}.html"), "w", encoding="utf-8") as f:
                            f.write(html)
                        log(f"  第 {idx+1}/{total_pages} 页完成", "ok")
                        return {"idx": idx, "html": html, "ok": True}
                    except Exception as e:
                        if attempt < 2:
                            log(f"  第 {idx+1} 页失败，重试... ({e})", "warn")
                            await asyncio.sleep(3)
                        else:
                            log(f"  第 {idx+1} 页最终失败: {e}", "err")
                            return {"idx": idx, "html": None, "ok": False}

        results = await asyncio.gather(*(gen_page(i) for i in range(total_pages)))

        sections = [None] * total_pages
        ok_count = 0
        for r in results:
            if r["ok"]:
                sections[r["idx"]] = r["html"]
                ok_count += 1

        if ok_count == 0:
            log("所有页面生成失败", "err")
            return False

        log(f"HTML 生成完成: {ok_count}/{total_pages} 页成功", "ok")

        # ── Phase 4: 组装完整 HTML ──
        log("Phase 4: 组装完整 HTML...")
        images_dir = os.path.join(work_dir, "images")
        images_b64 = {}
        if os.path.exists(images_dir):
            for f in os.listdir(images_dir):
                if f.endswith(('.jpg', '.jpeg', '.png')):
                    img_id = os.path.splitext(f)[0]
                    images_b64[img_id] = base64.b64encode(
                        open(os.path.join(images_dir, f), "rb").read()
                    ).decode()

        from backend.core.html_builder import assemble_html
        successful_sections = [s for s in sections if s is not None]
        full_html = assemble_html(successful_sections, images_b64, style)

        html_path = os.path.join(work_dir, "output.html")
        with open(html_path, "w", encoding="utf-8") as f:
            f.write(full_html)
        log("HTML 组装完成", "ok")

        # ── Phase 5: 导出 PPTX ──
        output_files = []

        if "image" in export_modes:
            log("Phase 5: 导出图片版 PPTX...")
            pptx_image_path = os.path.join(output_dir, "presentation_image.pptx")
            try:
                await asyncio.to_thread(_build_image_pptx_standalone, html_path, pptx_image_path)
                output_files.append(pptx_image_path)
                log(f"图片版导出完成: {pptx_image_path}", "ok")
            except Exception as e:
                log(f"图片版导出失败: {e}", "err")
                traceback.print_exc()

        if "editable" in export_modes:
            log("Phase 5: 导出可编辑版 PPTX...")
            pptx_edit_path = os.path.join(output_dir, "presentation_editable.pptx")
            try:
                from backend.core.layout_extractor import extract_layout
                layout_data = await asyncio.to_thread(extract_layout, html_path)

                from backend.core.pptx_builder import build_pptx
                await asyncio.to_thread(
                    build_pptx, layout_data, style, pptx_edit_path,
                    os.path.join(work_dir, "images")
                )
                output_files.append(pptx_edit_path)
                log(f"可编辑版导出完成: {pptx_edit_path}", "ok")
            except Exception as e:
                log(f"可编辑版导出失败: {e}", "err")
                traceback.print_exc()

        if output_files:
            log(f"全部完成！输出文件:", "ok")
            for f in output_files:
                log(f"  -> {f}")
            return True
        else:
            log("所有导出均失败", "err")
            return False

    finally:
        # 清理工作目录
        shutil.rmtree(work_dir, ignore_errors=True)


# ─── 独立的图片版 PPTX 构建（不依赖 routes.py）──────────────────────────────────

def _build_image_pptx_standalone(html_path: str, pptx_path: str):
    """用 Playwright 截图每个 section，构建图片版 PPTX。"""
    from playwright.sync_api import sync_playwright
    from pptx import Presentation
    from pptx.util import Emu

    screenshots = []
    tmp_dir = tempfile.mkdtemp()

    try:
        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page(viewport={"width": 1280, "height": 720})
            page.goto(f"file://{os.path.abspath(html_path)}")
            page.wait_for_load_state("networkidle")

            page_sections = page.query_selector_all("section")
            total = len(page_sections)
            log(f"  截图: 共 {total} 页")

            for i, section in enumerate(page_sections):
                img_path = os.path.join(tmp_dir, f"slide_{i+1:02d}.png")
                section.screenshot(path=img_path)
                screenshots.append(img_path)
                log(f"  截图: {i+1}/{total}")

            browser.close()

        # 构建 PPTX
        prs = Presentation()
        prs.slide_width = Emu(12192000)   # 13.333"
        prs.slide_height = Emu(6858000)   # 7.5"
        blank_layout = prs.slide_layouts[6]

        for img_path in screenshots:
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                img_path,
                left=Emu(0), top=Emu(0),
                width=prs.slide_width, height=prs.slide_height
            )

        prs.save(pptx_path)
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


# ─── CLI 入口 ──────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Prism CLI — AI-powered DOCX to PPT converter",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例:
  # 基本用法（使用环境变量中的 API key）
  python3 prism_cli.py --input slides.docx --output-dir ./output

  # 指定风格
  python3 prism_cli.py --input slides.docx --style xiaomi-light --output-dir ./output

  # 使用自定义 API
  python3 prism_cli.py --input slides.docx --api-key sk-xxx --base-url https://api.openai.com/v1 --model gpt-4o

  # 仅导出图片版
  python3 prism_cli.py --input slides.docx --export image

  # 查看可用风格
  python3 prism_cli.py --list-styles

  # 检查依赖
  python3 prism_cli.py --check-deps
""")

    parser.add_argument("--input", "-i", help="输入 .docx 文件路径")
    parser.add_argument("--output-dir", "-o", default="./output", help="输出目录 (默认: ./output)")
    parser.add_argument("--style", "-s", default="xiaomi-dark", help="风格名称 (默认: xiaomi-dark)")
    parser.add_argument("--export", "-e", default="both",
                        choices=["image", "editable", "both"],
                        help="导出模式 (默认: both)")
    parser.add_argument("--work-dir", help="中间产物目录（宿主模式使用）")
    parser.add_argument("--page", type=int, help="页码（用于单页 prompt 导出）")
    parser.add_argument("--intents-file", help="布局意图 JSON 文件路径")
    parser.add_argument("--prompt-output", help="将 prompt bundle 保存为 JSON 文件")
    parser.add_argument("--html-output", help="完整 HTML 输出路径")

    # API 配置（可选，覆盖环境变量和配置文件）
    api_group = parser.add_argument_group("API 配置（可选，覆盖默认配置）")
    api_group.add_argument("--api-key", help="API Key")
    api_group.add_argument("--base-url", help="API Base URL")
    api_group.add_argument("--model", help="模型名称")
    api_group.add_argument("--timeout", type=int, help="超时时间（秒）")

    # 工具命令
    parser.add_argument("--list-styles", action="store_true", help="列出所有可用风格")
    parser.add_argument("--check-deps", action="store_true", help="检查依赖安装情况")
    parser.add_argument("--parse-docx", action="store_true", help="解析 DOCX 到 work_dir/pages.json")
    parser.add_argument("--emit-layout-plan", action="store_true", help="导出布局规划 prompt bundle")
    parser.add_argument("--emit-page-prompt", action="store_true", help="导出单页 HTML 生成 prompt bundle")
    parser.add_argument("--assemble-html", action="store_true", help="将 sections/*.html 组装为完整 HTML")
    parser.add_argument("--export-from-workdir", action="store_true", help="从 work_dir 导出 PPTX")

    args = parser.parse_args()

    # ── 工具命令 ──
    if args.list_styles:
        styles = list_styles()
        print("\n可用风格:")
        for s in styles:
            mixed_tag = " [混合明暗]" if s["mixed"] else ""
            print(f"  {s['id']:30s} {s['name']}{mixed_tag}")
        return

    if args.check_deps:
        print("\n依赖检查:")
        deps = check_dependencies()
        all_ok = True
        for name, ok in deps.items():
            status = "OK" if ok else "MISSING"
            if not ok:
                all_ok = False
            print(f"  {name:25s} {status}")
        if not all_ok:
            print("\n部分依赖缺失，请运行 setup.sh 安装")
            sys.exit(1)
        else:
            print("\n所有依赖已就绪")
        return

    if args.parse_docx:
        if not args.input:
            parser.error("--parse-docx 需要指定 --input")
        if not args.work_dir:
            parser.error("--parse-docx 需要指定 --work-dir")
        pages_data = prepare_docx_workdir(args.input, args.work_dir)
        print(f"\n解析完成：共 {pages_data.get('total_pages', 0)} 页")
        print(f"pages.json: {Path(args.work_dir) / 'pages.json'}")
        return

    if args.emit_layout_plan:
        if not args.work_dir:
            parser.error("--emit-layout-plan 需要指定 --work-dir")
        bundle = emit_layout_plan_prompt_bundle(args.work_dir, args.style, args.prompt_output)
        print(json.dumps(bundle, ensure_ascii=False, indent=2))
        return

    if args.emit_page_prompt:
        if not args.work_dir:
            parser.error("--emit-page-prompt 需要指定 --work-dir")
        if not args.page:
            parser.error("--emit-page-prompt 需要指定 --page")
        bundle = emit_page_prompt_bundle(
            args.work_dir,
            args.style,
            args.page,
            args.intents_file,
            args.prompt_output,
        )
        print(json.dumps(bundle, ensure_ascii=False, indent=2))
        return

    if args.assemble_html:
        if not args.work_dir:
            parser.error("--assemble-html 需要指定 --work-dir")
        html_path = assemble_html_from_workdir(args.work_dir, args.style, args.html_output)
        print(f"\nHTML 组装完成：{html_path}")
        return

    if args.export_from_workdir:
        if not args.work_dir:
            parser.error("--export-from-workdir 需要指定 --work-dir")
        export_modes = {
            "image": ["image"],
            "editable": ["editable"],
            "both": ["image", "editable"],
        }[args.export]
        files = export_from_workdir(args.work_dir, args.output_dir, args.style, export_modes)
        print("\n导出完成：")
        for file_path in files:
            print(f"  {file_path}")
        return

    # ── 主流程 ──
    if not args.input:
        parser.error("需要指定 --input 参数")

    if not os.path.exists(args.input):
        print(f"错误: 文件不存在: {args.input}")
        sys.exit(1)

    if not args.input.lower().endswith('.docx'):
        print(f"错误: 请提供 .docx 文件")
        sys.exit(1)

    # 设置 API 配置覆盖
    if args.api_key or args.base_url or args.model:
        from backend.core.llm_client import set_runtime_api_config
        set_runtime_api_config(
            api_key=args.api_key,
            base_url=args.base_url,
            model=args.model,
            timeout=args.timeout,
        )

    # 确定导出模式
    export_modes = {
        "image": ["image"],
        "editable": ["editable"],
        "both": ["image", "editable"],
    }[args.export]

    print(f"\n{'='*60}")
    print(f"  Prism CLI — DOCX to PPT")
    print(f"{'='*60}")
    print(f"  输入: {args.input}")
    print(f"  风格: {args.style}")
    print(f"  导出: {', '.join(export_modes)}")
    print(f"  输出: {os.path.abspath(args.output_dir)}")
    print(f"{'='*60}\n")

    success = asyncio.run(
        run_cli_pipeline(
            docx_path=args.input,
            style_name=args.style,
            output_dir=args.output_dir,
            export_modes=export_modes,
        )
    )

    sys.exit(0 if success else 1)


if __name__ == "__main__":
    main()
