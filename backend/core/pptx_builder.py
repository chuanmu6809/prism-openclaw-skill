from __future__ import annotations
"""
pptx_builder.py
从 layout_data.json 重建 PPT。
"""
import re
from pathlib import Path
from io import BytesIO

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
from PIL import Image

from backend.core.paths import PROJECT_ROOT

_ROOT = PROJECT_ROOT
_IMAGES_DIR = _ROOT / 'output' / 'images'  # default fallback

# ─── 颜色解析 ──────────────────────────────────────────────────────────────────

_HEX_ALPHA_RE = re.compile(r'^#([0-9A-Fa-f]{6}):?([\d.]+)?$')

def _parse_color(color_str: str) -> tuple[RGBColor, float] | None:
    """
    解析颜色字符串，返回 (RGBColor, alpha) 或 None。
    支持格式：
      #RRGGBB        → alpha=1.0
      #RRGGBB:0.70   → alpha=0.70
      transparent     → None
    """
    if not color_str or color_str == 'transparent':
        return None
    m = _HEX_ALPHA_RE.match(color_str)
    if m:
        rgb = RGBColor.from_string(m.group(1))
        alpha = float(m.group(2)) if m.group(2) else 1.0
        return rgb, alpha
    return None

# ─── 对齐映射 ──────────────────────────────────────────────────────────────────

_ALIGN_MAP = {
    'center': PP_ALIGN.CENTER,
    'left': PP_ALIGN.LEFT,
    'start': PP_ALIGN.LEFT,
    'right': PP_ALIGN.RIGHT,
    'end': PP_ALIGN.RIGHT,
    'justify': PP_ALIGN.JUSTIFY,
}

# ─── XML 工具 ──────────────────────────────────────────────────────────────────

_NSMAP = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
}

def _set_transparency(fill, alpha: float):
    """给 solidFill 设置透明度（alpha < 1.0 时）。"""
    if alpha >= 1.0:
        return
    # python-pptx 不直接支持透明度，需要操作 XML
    # 通过 fill 的 xml element 找到 srgbClr 节点
    try:
        fill_xml = fill._fill if hasattr(fill, '_fill') else fill._element
        srgb = fill_xml.find(f'.//{{{_NSMAP["a"]}}}srgbClr')
        if srgb is not None:
            # alpha 属性值：0-100000（万分之一）
            alpha_el = etree.SubElement(srgb, f'{{{_NSMAP["a"]}}}alpha')
            alpha_el.set('val', str(int(alpha * 100000)))
    except (AttributeError, TypeError):
        pass  # 无法设置透明度时跳过


def _set_image_opacity(pic, opacity: float):
    """给图片设置透明度（通过 alphaModFix）。"""
    if opacity >= 1.0:
        return
    try:
        blip_fill = pic._element.find(f'.//{{{_NSMAP["a"]}}}blip')
        if blip_fill is not None:
            alpha_mod = etree.SubElement(blip_fill, f'{{{_NSMAP["a"]}}}alphaModFix')
            alpha_mod.set('amt', str(int(opacity * 100000)))
    except (AttributeError, TypeError):
        pass

def _set_corner_radius(shape, radius_px: int):
    """给形状设置圆角。"""
    if radius_px <= 0:
        return
    # 圆角单位是 EMU
    radius_emu = radius_px * 9525
    sp = shape._element
    prstGeom = sp.find(f'.//{{{_NSMAP["a"]}}}prstGeom')
    if prstGeom is not None:
        avLst = prstGeom.find(f'{{{_NSMAP["a"]}}}avLst')
        if avLst is None:
            avLst = etree.SubElement(prstGeom, f'{{{_NSMAP["a"]}}}avLst')
        # 圆角值需要计算：val = radius / min(width, height) * 50000
        w = shape.width
        h = shape.height
        min_dim = min(w, h) if min(w, h) > 0 else 1
        val = int(radius_emu / min_dim * 50000)
        val = min(val, 50000)  # 最大 50000
        gd = etree.SubElement(avLst, f'{{{_NSMAP["a"]}}}gd')
        gd.set('name', 'adj')
        gd.set('fmla', f'val {val}')


def _set_gradient(fill, gradient: dict):
    """设置线性渐变填充。"""
    fill.gradient()
    fill.gradient_angle = gradient.get('angle', 180)
    
    c1 = _parse_color(gradient.get('color1', ''))
    c2 = _parse_color(gradient.get('color2', ''))
    
    stops = fill.gradient_stops
    if c1:
        stops[0].color.rgb = c1[0]
        stops[0].position = 0.0
    if c2:
        stops[1].color.rgb = c2[0]
        stops[1].position = 1.0
    
    # 处理渐变色的透明度
    if c1 and c1[1] < 1.0:
        _set_gradient_stop_alpha(stops[0], c1[1])
    if c2 and c2[1] < 1.0:
        _set_gradient_stop_alpha(stops[1], c2[1])


def _set_gradient_stop_alpha(stop, alpha: float):
    """给渐变停靠点设置透明度。"""
    # 通过 XML 操作
    gs_el = stop._element
    srgb = gs_el.find(f'{{{_NSMAP["a"]}}}srgbClr')
    if srgb is not None:
        alpha_el = etree.SubElement(srgb, f'{{{_NSMAP["a"]}}}alpha')
        alpha_el.set('val', str(int(alpha * 100000)))


# ─── 图片裁切计算 ──────────────────────────────────────────────────────────────

def _calc_crop(img_path: Path, container_w_emu: int, container_h_emu: int,
               object_fit: str) -> dict:
    """
    计算 object-fit: cover 的裁切参数。
    返回 {crop_left, crop_right, crop_top, crop_bottom}，值为 0.0-1.0。
    """
    if object_fit != 'cover':
        return {'crop_left': 0, 'crop_right': 0, 'crop_top': 0, 'crop_bottom': 0}

    try:
        with Image.open(img_path) as img:
            img_w, img_h = img.size
    except Exception:
        return {'crop_left': 0, 'crop_right': 0, 'crop_top': 0, 'crop_bottom': 0}

    # CSS object-fit: cover 的计算逻辑
    container_ratio = container_w_emu / container_h_emu if container_h_emu else 1
    img_ratio = img_w / img_h if img_h else 1

    if img_ratio > container_ratio:
        # 图片更宽 → 上下铺满，左右裁切
        scale = container_h_emu / img_h
        scaled_w = img_w * scale
        overflow = (scaled_w - container_w_emu) / scaled_w
        return {
            'crop_left': overflow / 2,
            'crop_right': overflow / 2,
            'crop_top': 0,
            'crop_bottom': 0,
        }
    else:
        # 图片更高 → 左右铺满，上下裁切
        scale = container_w_emu / img_w
        scaled_h = img_h * scale
        overflow = (scaled_h - container_h_emu) / scaled_h
        return {
            'crop_left': 0,
            'crop_right': 0,
            'crop_top': overflow / 2,
            'crop_bottom': overflow / 2,
        }


# ─── 元素排序 ──────────────────────────────────────────────────────────────────

def _sort_elements(elements: list) -> list:
    """按 z-order 排序：背景图 → 装饰shape(大→小) → 内容图(大→小) → 文字。"""
    bg_imgs = []    # 全页背景图
    shapes = []     # 装饰 shape 截图
    imgs = []       # 内容图截图
    texts = []      # 文字

    for el in elements:
        t = el.get('type')
        area = el.get('w_emu', 0) * el.get('h_emu', 0)

        if t == 'container_image':
            if el.get('_is_image'):
                # 图片截图 → 放在 shapes 上方
                imgs.append((area, el))
            else:
                # shape 截图 → 底层
                shapes.append((area, el))
        elif t == 'image':
            # 面积接近全页的视为背景图
            full_area = 12192000 * 6858000
            if area > full_area * 0.8:
                bg_imgs.append((area, el))
            else:
                imgs.append((area, el))
        elif t in ('shape',):
            shapes.append((area, el))
        elif t == 'text':
            texts.append((area, el))

    # 每层内部：面积大的先添加（下层）
    result = []
    result.extend(el for _, el in sorted(bg_imgs, key=lambda x: -x[0]))
    result.extend(el for _, el in sorted(shapes, key=lambda x: -x[0]))
    result.extend(el for _, el in sorted(imgs, key=lambda x: -x[0]))
    result.extend(el for _, el in sorted(texts, key=lambda x: -x[0]))

    return result


# ─── 主入口 ────────────────────────────────────────────────────────────────────

def build_pptx(layout_data: list, style: dict, output_path: str, images_dir: str = None):
    """
    从 layout_data 生成 PPTX。
    """
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)

    blank_layout = prs.slide_layouts[6]  # 空白版式
    resolved_images_dir = Path(images_dir) if images_dir else _IMAGES_DIR

    for page_data in layout_data:
        page_num = page_data['page']
        print(f'  → 重建第 {page_num} 页...')

        slide = prs.slides.add_slide(blank_layout)

        # ── 应用页面背景 ──
        bg_screenshot = page_data.get('bg_screenshot', '')
        if bg_screenshot and Path(bg_screenshot).exists():
            # 背景合成截图 → 全页底图
            slide.shapes.add_picture(
                bg_screenshot,
                Emu(0), Emu(0),
                prs.slide_width, prs.slide_height
            )
        else:
            # 回退：纯色背景
            page_bg = page_data.get('bg_color', '')
            parsed_bg = _parse_color(page_bg) if page_bg else None
            if parsed_bg:
                bg_fill = slide.background.fill
                bg_fill.solid()
                bg_fill.fore_color.rgb = parsed_bg[0]

        # ── 按 z-order 排序元素 ──
        elements = _sort_elements(page_data.get('elements', []))

        for el in elements:
            el_type = el.get('type')

            if el_type == 'container_image':
                _add_container_image(slide, el)
            elif el_type == 'text':
                _add_text(slide, el)

        print(f'    ✓ {len(elements)} 个元素')

    prs.save(output_path)
    print(f'\n  PPT 已保存：{output_path}')


# ─── 元素添加函数 ──────────────────────────────────────────────────────────────

def _add_container_image(slide, el: dict):
    """添加容器截图（作为图片插入，文字单独叠加）。"""
    img_path = el.get('screenshot_path', '')
    if not img_path or not Path(img_path).exists():
        print(f'    [警告] 容器截图不存在：{img_path}')
        return

    left = Emu(el['x_emu'])
    top = Emu(el['y_emu'])
    width = Emu(el['w_emu'])
    height = Emu(el['h_emu'])

    pic = slide.shapes.add_picture(
        str(img_path), left, top, width, height
    )

    # 圆角
    radius_px = el.get('border_radius_px', 0)
    if radius_px > 0:
        sp = pic._element
        spPr = sp.find(f'.//{{{_NSMAP["a"]}}}prstGeom')
        if spPr is not None:
            spPr.set('prst', 'roundRect')
            avLst = spPr.find(f'{{{_NSMAP["a"]}}}avLst')
            if avLst is None:
                avLst = etree.SubElement(spPr, f'{{{_NSMAP["a"]}}}avLst')
            radius_emu = radius_px * 9525
            min_dim = min(pic.width, pic.height) if min(pic.width, pic.height) > 0 else 1
            val = int(radius_emu / min_dim * 50000)
            val = min(val, 50000)
            gd = etree.SubElement(avLst, f'{{{_NSMAP["a"]}}}gd')
            gd.set('name', 'adj')
            gd.set('fmla', f'val {val}')

    # 透明度
    opacity = el.get('opacity', 1.0)
    if opacity < 1.0:
        _set_image_opacity(pic, opacity)

def _add_image(slide, el: dict, images_dir: Path = None):
    """添加图片元素。"""
    img_id = el.get('img_id', 'UNKNOWN')
    resolved = images_dir or _IMAGES_DIR
    img_path = resolved / f'{img_id}.jpg'

    if not img_path.exists():
        print(f'    [警告] 图片文件不存在：{img_path}')
        return

    left = Emu(el['x_emu'])
    top = Emu(el['y_emu'])
    width = Emu(el['w_emu'])
    height = Emu(el['h_emu'])

    pic = slide.shapes.add_picture(
        str(img_path), left, top, width, height
    )

    # 计算裁切参数（object-fit: cover）
    object_fit = el.get('object_fit', 'cover')
    crop = _calc_crop(img_path, el['w_emu'], el['h_emu'], object_fit)
    pic.crop_left = crop['crop_left']
    pic.crop_right = crop['crop_right']
    pic.crop_top = crop['crop_top']
    pic.crop_bottom = crop['crop_bottom']

    # 圆角（通过 XML 修改几何形状）
    radius_px = el.get('border_radius_px', 0)
    if radius_px > 0:
        sp = pic._element
        spPr = sp.find(f'.//{{{_NSMAP["a"]}}}prstGeom')
        if spPr is not None:
            spPr.set('prst', 'roundRect')
            avLst = spPr.find(f'{{{_NSMAP["a"]}}}avLst')
            if avLst is None:
                avLst = etree.SubElement(spPr, f'{{{_NSMAP["a"]}}}avLst')
            radius_emu = radius_px * 9525
            min_dim = min(pic.width, pic.height) if min(pic.width, pic.height) > 0 else 1
            val = int(radius_emu / min_dim * 50000)
            val = min(val, 50000)
            gd = etree.SubElement(avLst, f'{{{_NSMAP["a"]}}}gd')
            gd.set('name', 'adj')
            gd.set('fmla', f'val {val}')


def _add_shape(slide, el: dict):
    """添加形状元素（矩形/装饰条）。"""
    # 跳过完全透明且无边框的 shape（无渐变、背景 alpha=0、无 border）
    gradient = el.get('gradient')
    border = el.get('border')
    if not gradient and not border:
        bg_color = el.get('bg_color', '')
        parsed = _parse_color(bg_color)
        if parsed and parsed[1] == 0:
            return  # 完全透明且无边框，跳过

    left = Emu(el['x_emu'])
    top = Emu(el['y_emu'])
    width = Emu(max(el['w_emu'], 1))
    height = Emu(max(el['h_emu'], 1))

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if el.get('border_radius_px', 0) > 0 else MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )

    # 填充
    gradient = el.get('gradient')
    if gradient:
        _set_gradient(shape.fill, gradient)
    else:
        bg_color = el.get('bg_color', '')
        parsed = _parse_color(bg_color)
        if parsed:
            shape.fill.solid()
            shape.fill.fore_color.rgb = parsed[0]
            if parsed[1] < 1.0:
                _set_transparency(shape.fill, parsed[1])
        else:
            shape.fill.background()  # 透明

    # 边框 / 描边
    border = el.get('border')
    if border:
        border_color = _parse_color(border.get('color', ''))
        if border_color:
            shape.line.color.rgb = border_color[0]
            shape.line.width = Pt(border.get('width_pt', 1))
            if border_color[1] < 1.0:
                # 边框透明度通过 XML
                ln = shape._element.find(f'.//{{{_NSMAP["a"]}}}ln')
                if ln is not None:
                    srgb = ln.find(f'.//{{{_NSMAP["a"]}}}srgbClr')
                    if srgb is not None:
                        alpha_el = etree.SubElement(srgb, f'{{{_NSMAP["a"]}}}alpha')
                        alpha_el.set('val', str(int(border_color[1] * 100000)))
    else:
        shape.line.fill.background()  # 无边框

    # 圆角
    if el.get('border_radius_px', 0) > 0:
        _set_corner_radius(shape, el['border_radius_px'])


def _add_text(slide, el: dict):
    """添加文字元素。"""
    left = Emu(el['x_emu'])
    top = Emu(el['y_emu'])
    width = Emu(max(el['w_emu'], 1))
    height = Emu(max(el['h_emu'], 1))

    txBox = slide.shapes.add_textbox(left, top, width, height)
    # 显式去除边框
    txBox.line.fill.background()
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None  # 不自动缩放

    # 将内边距设为 0，避免文字位置偏移
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)

    # 纵向对齐
    v_align = el.get('vertical_align', 'top')
    if v_align == 'middle':
        from pptx.enum.text import MSO_ANCHOR
        tf.word_wrap = True
        # 通过 XML 设置 bodyPr 的 anchor 属性
        bodyPr = txBox._element.find(f'.//{{{_NSMAP["a"]}}}bodyPr')
        if bodyPr is not None:
            bodyPr.set('anchor', 'ctr')

    text_align = el.get('text_align', 'left')
    runs_data = el.get('runs')

    if runs_data:
        # 多 run 模式：混合字体段落
        p = tf.paragraphs[0]
        p.alignment = _ALIGN_MAP.get(text_align, PP_ALIGN.LEFT)

        for run_info in runs_data:
            run_text = run_info.get('text', '')
            if not run_text:
                continue

            # 处理含换行的 run
            lines = run_text.split('\n')
            for li, line in enumerate(lines):
                if li > 0:
                    p = tf.add_paragraph()
                    p.alignment = _ALIGN_MAP.get(text_align, PP_ALIGN.LEFT)

                if not line:
                    continue

                run = p.add_run()
                run.text = line

                # 字号：优先用 run 自己的，否则用元素级的
                font_size_pt = run_info.get('font_size_pt', el.get('font_size_pt', 15))
                run.font.size = Pt(font_size_pt)

                # 字体
                font_family = run_info.get('font_family', el.get('font_family', 'MiSans-Regular'))
                run.font.name = font_family

                # 颜色
                color_str = run_info.get('color', el.get('color', '#FFFFFF'))
                parsed_color = _parse_color(color_str)
                if parsed_color:
                    run.font.color.rgb = parsed_color[0]
                    alpha = parsed_color[1]
                    if alpha < 1.0:
                        rPr = run._r.get_or_add_rPr()
                        solidFill = rPr.find(f'{{{_NSMAP["a"]}}}solidFill')
                        if solidFill is not None:
                            srgbClr = solidFill.find(f'{{{_NSMAP["a"]}}}srgbClr')
                            if srgbClr is not None:
                                alpha_el = etree.SubElement(srgbClr, f'{{{_NSMAP["a"]}}}alpha')
                                alpha_el.set('val', str(int(alpha * 100000)))

                # 粗体判断
                font_weight = run_info.get('font_weight', el.get('font_weight', '400'))
                if font_weight in ('700', '600', 'bold') or 'Bold' in font_family or 'Semibold' in font_family:
                    run.font.bold = True
    else:
        # 单一字体模式（原有逻辑）
        text = el.get('text', '')
        color_str = el.get('color', '#FFFFFF')
        font_size_pt = el.get('font_size_pt', 15)
        font_family = el.get('font_family', 'MiSans-Regular')
        font_weight = el.get('font_weight', '400')

        p = tf.paragraphs[0]
        p.alignment = _ALIGN_MAP.get(text_align, PP_ALIGN.LEFT)

        # 处理多行文字（innerText 可能包含换行）
        lines = text.split('\n')
        for i, line in enumerate(lines):
            if i > 0:
                p = tf.add_paragraph()
                p.alignment = _ALIGN_MAP.get(text_align, PP_ALIGN.LEFT)

            run = p.add_run()
            run.text = line
            run.font.size = Pt(font_size_pt)
            run.font.name = font_family

            parsed_color = _parse_color(color_str)
            if parsed_color:
                run.font.color.rgb = parsed_color[0]
                # 文字颜色透明度（alpha < 1.0 时）
                alpha = parsed_color[1]
                if alpha < 1.0:
                    rPr = run._r.get_or_add_rPr()
                    solidFill = rPr.find(f'{{{_NSMAP["a"]}}}solidFill')
                    if solidFill is not None:
                        srgbClr = solidFill.find(f'{{{_NSMAP["a"]}}}srgbClr')
                        if srgbClr is not None:
                            alpha_el = etree.SubElement(srgbClr, f'{{{_NSMAP["a"]}}}alpha')
                            alpha_el.set('val', str(int(alpha * 100000)))

            # 粗体判断
            if font_weight in ('700', '600', 'bold') or 'Bold' in font_family or 'Semibold' in font_family:
                run.font.bold = True
